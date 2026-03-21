import os
import re
from pathlib import Path
from typing import Dict, List

import fitz
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openai import OpenAI
from pptx import Presentation

from .db import insert_chunks, source_already_ingested


def infer_chapter(filename: str) -> str:
    name = Path(filename).stem
    match = re.search(r"chapter\s*[_-]?(\d+)", name, re.IGNORECASE)
    if match:
        return f"Chapter {match.group(1)}"
    return name


def extract_pdf_pages(file_path: Path) -> List[Dict]:
    pages = []
    doc = fitz.open(file_path)
    try:
        for i, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            if text:
                pages.append({"page_number": i, "text": text})
    finally:
        doc.close()
    return pages


def extract_pptx_pages(file_path: Path) -> List[Dict]:
    pages = []
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                parts.append(notes)
        text = "\n".join(parts).strip()
        if text:
            pages.append({"page_number": i, "text": text})
    return pages


def extract_docx_pages(file_path: Path) -> List[Dict]:
    doc = Document(file_path)
    text = "\n".join(p.text for p in doc.paragraphs).strip()
    if not text:
        return []
    return [{"page_number": 1, "text": text}]


def discover_supported_files(folder_path: str) -> List[Path]:
    base = Path(folder_path)
    files = []
    for ext in ("*.pdf", "*.pptx", "*.docx"):
        files.extend(base.rglob(ext))
    return files


def ingest_folder(path: str) -> Dict[str, int]:
    files = discover_supported_files(path)
    splitter = RecursiveCharacterTextSplitter(chunk_size=400, chunk_overlap=60)
    client = OpenAI()

    files_processed = 0
    chunks_stored = 0

    for file_path in files:
        source = file_path.name
        if source_already_ingested(source):
            continue

        ext = file_path.suffix.lower()
        if ext == ".pdf":
            pages = extract_pdf_pages(file_path)
        elif ext == ".pptx":
            pages = extract_pptx_pages(file_path)
        elif ext == ".docx":
            pages = extract_docx_pages(file_path)
        else:
            continue

        chunk_rows = []
        chapter = infer_chapter(source)

        for page in pages:
            chunks = splitter.split_text(page["text"])
            for chunk in chunks:
                clean_chunk = chunk.strip()
                if not clean_chunk:
                    continue
                chunk_rows.append(
                    {
                        "source": source,
                        "chapter": chapter,
                        "page_number": page["page_number"],
                        "content": clean_chunk,
                    }
                )

        if not chunk_rows:
            files_processed += 1
            continue

        embeddings_response = client.embeddings.create(
            model="text-embedding-3-small",
            input=[row["content"] for row in chunk_rows],
        )
        embeddings = [item.embedding for item in embeddings_response.data]

        insert_rows = []
        for row, embedding in zip(chunk_rows, embeddings):
            row_with_embedding = dict(row)
            row_with_embedding["embedding"] = embedding
            insert_rows.append(row_with_embedding)

        chunks_stored += insert_chunks(insert_rows)
        files_processed += 1

    return {"files_processed": files_processed, "chunks_stored": chunks_stored}
